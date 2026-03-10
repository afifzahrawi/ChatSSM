"""
preprocess.py  –  ChatSSM Offline Preprocessor
================================================
Run ONCE per source (or whenever a source changes) to convert PDFs and CSVs
into clean, structured CSV files ready for the Streamlit app to index.

Usage
-----
    python preprocess.py                         # all discovered sources
    python preprocess.py --key pn_1_2023         # one source by key
    python preprocess.py --key pn_1_2023 --force # force re-process
    python preprocess.py --list                  # show status of all sources
    python preprocess.py --scan                  # preview discovered files (dry-run)

═══════════════════════════════════════════════════════════
HOW TO ADD DOCUMENTS  (no JSON editing required)
═══════════════════════════════════════════════════════════

Just drop your PDF into the right folder:

    knowledge_base/sources/<Category>/<doc_type>/<filename>.pdf

Examples:
    knowledge_base/sources/Legislations/act/companies_act_2016.pdf
    knowledge_base/sources/Practice Notes/general/pn_3_2018.pdf
    knowledge_base/sources/FAQ/faq/faq_incorporation.pdf
    knowledge_base/sources/Guidelines/general/mccg_2021.pdf
    knowledge_base/sources/Circular/general/circular_1_2024.pdf

Valid <Category> folder names:
    Legislations | Practice Notes | Practice Directives |
    Guidelines   | Circular       | FAQ | Forms

Valid <doc_type> folder names:
    act      – Malaysian Acts (PART → Division → Section structure)
    general  – Practice Notes, Guidelines, Circulars (headings + paragraphs)
    faq      – FAQ documents (Q: / A: pairs)

The key is derived from the filename:
    "PN No. 3-2018 Extension of Time.pdf"  →  key = "pn_no_3_2018_extension_of_time"

The display name is the filename without extension (can be overridden in JSON).

═══════════════════════════════════════════════════════════
OPTIONAL: knowledge_sources.json overrides
═══════════════════════════════════════════════════════════

knowledge_sources.json is now OPTIONAL. Use it only when you need to:
  • Give a document a custom display name  (name field)
  • Point to a remote URL instead of a local file  (url field)
  • Disable a file without deleting it  (enabled: false)
  • Add a source that lives outside the sources/ folder

Files discovered from folders take precedence; JSON entries can override
any auto-discovered field or add sources that folders don't cover.

Output
------
    knowledge_base/processed/<key>.csv
    Columns: source_key, source_name, part, division, subdivision,
             section, section_title, content
"""

from __future__ import annotations

import argparse
import csv
import io
import json
import logging
import os
import re
import sys
import shutil
import tempfile
from dataclasses import dataclass, field
from typing import Dict, List, Optional
from pypdf import PdfReader
from concurrent.futures import ThreadPoolExecutor, as_completed

import pdfplumber
import requests

# ─── Logging ─────────────────────────────────────────────────────────────────

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger("preprocess")

# ─── Paths ────────────────────────────────────────────────────────────────────

SOURCES_CONFIG   = "knowledge_sources.json"   # optional overrides file
SOURCES_DIR      = os.path.join("knowledge_base", "sources")   # auto-discovery root
PROCESSED_DIR    = os.path.join("knowledge_base", "processed")
os.makedirs(PROCESSED_DIR, exist_ok=True)
os.makedirs(SOURCES_DIR,   exist_ok=True)

# Minimum character length to keep a section (filters out headers, TOC noise)
MIN_CONTENT_CHARS = 60

# ─── OCR / Vision settings ───────────────────────────────────────────────────
OCR_MODEL          = os.environ.get("CHATSSM_OCR_MODEL", "glm-ocr")
OCR_TIMEOUT        = int(os.environ.get("CHATSSM_OCR_TIMEOUT", "600"))   # seconds per page/image
MIN_TEXT_PER_PAGE  = 100   # chars; pages below this are treated as scanned and sent to OCR

# ─── GPU offload settings (Ollama) ──────────────────────────────────────────
# -1 = all layers on GPU; 0 = CPU only; N = offload N layers
OCR_NUM_GPU    = int(os.environ.get("CHATSSM_OCR_NUM_GPU",   "-1"))
CHUNK_NUM_GPU  = int(os.environ.get("CHATSSM_CHUNK_NUM_GPU", "-1"))
OCR_WORKERS    = int(os.environ.get("CHATSSM_OCR_WORKERS",    "4"))

# ─── CSV output columns ───────────────────────────────────────────────────────

CSV_COLUMNS = [
    "source_key",
    "source_name",
    "part",
    "division",
    "subdivision",
    "section",
    "section_title",
    "content",
]

# ─── Regex patterns for Malaysian legal document structure ────────────────────
#
# CRITICAL DESIGN NOTES
# ─────────────────────
# 1. _PART_RE must NOT use re.IGNORECASE.
#    With IGNORECASE the pattern matches "part in the management..." because
#    [IVXLCDM] matches lowercase 'i', and [A-Z]? with IGNORECASE matches 'n',
#    so "part in" is parsed as "PART I+n". Malaysian Acts always write PART in
#    uppercase (or "Part" in title-case for LLP Act). The lookahead (?!\s+[a-z])
#    is a belt-and-braces guard: a real PART header is never followed by a
#    lowercase word ("in", "of", "for", ...).
#
# 2. _DIV_RE and _SUBDIV_RE keep IGNORECASE because "Division" / "Subdivision"
#    are proper words that only appear as structural headers — no false-positive
#    risk with lowercase.
#
# 3. _SECTION_RE intentionally keeps IGNORECASE so "Section 14." and "14."
#    both match.

# PART I, PART IX, PART IVA, PART 1, Part II (title-case used in LLP Act)
# NOT "part in …", NOT "part of …"  (the (?!\s+[a-z]) lookahead blocks these)
_PART_RE = re.compile(
    r"^(PART|Part)\s+([IVXLCDM]+[A-Za-z]?|\d+[A-Za-z]?)\b(?!\s+[a-z])"
)
# Division 1, Division 1A
_DIV_RE = re.compile(r"^Division\s+(\d+[A-Z]?)\b", re.IGNORECASE)
# Subdivision 1, Subdivision 1A
_SUBDIV_RE = re.compile(r"^Subdivision\s+(\d+[A-Z]?)\b", re.IGNORECASE)
# Section number: "14.", "14A.", "196.", optionally preceded by "Section "
# group(1) = number, group(2) = rest of line (may be the inline section title)
_SECTION_RE = re.compile(r"^(?:Section\s+)?(\d+[A-Z]?)\.\s*(.*)", re.IGNORECASE)
# Table-of-contents entry: ends with dots then a page number
_TOC_RE = re.compile(r"\.{3,}\s*\d+\s*$")
# Header/footer noise: standalone page numbers, "P.U. (A) 234" gazette refs
_PAGE_RE = re.compile(r"^\s*\d+\s*$")
_GAZETTE_RE = re.compile(r"^P\.U\.\s*\([AB]\)\s*\d+", re.IGNORECASE)
# "Laws of Malaysia ACT 743" / "Laws of Malaysia" — running header in all Acts
_LAWS_HDR_RE = re.compile(r"^Laws of Malaysia\b", re.IGNORECASE)
# "NOTE-The Companies Act 1965 [Act 125] has been repealed…" — editorial notes
_NOTE_RE = re.compile(r"^NOTE[-–]", re.IGNORECASE)

# ── Marginal note detector ────────────────────────────────────────────────────
# Malaysian Acts place a short descriptive heading (marginal note) in the right
# margin next to each section.  pdfplumber reads these as the LAST LINE of the
# preceding section's content block instead of the title of the current section.
#
# A line is treated as a marginal note when ALL of the following hold:
#   • Starts with a capital letter then a lowercase letter  (Title Case, not ALL CAPS)
#   • Length ≤ 80 characters  (headings are brief)
#   • No sentence-ending punctuation (.  !  ?)  — it is a heading, not a sentence
#   • No digits  — section references like "(1)" appear in content, not headings
#   • No parentheses — content lines start with "(1)", marginal notes don't
#   • Does not start with a common content word ("the", "a", "an", "in", etc.)
#
# The detector is intentionally conservative: a false negative (missing a
# marginal note) is much less harmful than a false positive (stripping real
# legal content).
_MARGINAL_NOTE_RE = re.compile(
    r"^[A-Z][a-z]"           # Title Case: capital + immediately lowercase
    r"[^.!?\d\(\);]{4,78}"   # 4-78 more chars; no sentence-ending punctuation,
                              # no digits, no parentheses, no semicolons
    r"$"                      # end of line — must be a standalone line
)
_CONTENT_STARTERS = re.compile(
    r"^(the|a|an|in|of|on|or|and|any|no|if|for|to|by|be|at|as|is|are|was|"
    r"were|has|have|had|its|his|her|their|shall|must|may|where|when|who|which"
    r")\b",
    re.IGNORECASE,
)


def _is_marginal_note(line: str) -> bool:
    """
    Return True if *line* looks like a Malaysian Act marginal note heading.
    See _MARGINAL_NOTE_RE docstring for the full set of criteria.
    """
    if not _MARGINAL_NOTE_RE.match(line):
        return False
    if _CONTENT_STARTERS.match(line):
        return False
    # Reject lines that look like continuation of a subsection reference
    # e.g. "(a) with the leave of ..." — already blocked by no-parentheses rule
    # Extra guard: must not contain a colon followed by text (structural header)
    # Allow colons for titles like "Compliance officer: responsibilities"
    return True


# ─── Data class ───────────────────────────────────────────────────────────────

@dataclass
class SourceEntry:
    key:        str
    name:       str
    category:   str
    type:       str           # "pdf" or "csv"
    doc_type:   str           # "act" | "general" | "faq"  (controls which parser)
    enabled:    bool
    relates_to_acts:  List[str] = field(default_factory=list)
    url:        Optional[str] = None
    local_path: Optional[str] = None

    @property
    def output_path(self) -> str:
        return os.path.join(PROCESSED_DIR, f"{self.key}.csv")

    @property
    def is_processed(self) -> bool:
        return os.path.exists(self.output_path)


# ─── Valid folder names ───────────────────────────────────────────────────────

# Category folder names must match these exactly (case-sensitive on Linux)
VALID_CATEGORIES: List[str] = [
    "Legislations",
    "Practice Notes",
    "Practice Directives",
    "Guidelines",
    "Circular",
    "FAQ",
    "Forms",
]

# doc_type folder names
VALID_DOC_TYPES: List[str] = ["act", "general", "faq", "gazette", "slide", "others"]


def _filename_to_key(filename: str) -> str:
    """
    Convert a filename to a clean, unique source key.

    "PN No. 3-2018 Extension of Time.pdf"
    → "pn_no_3_2018_extension_of_time"

    Rules:
    - Strip extension
    - Lowercase
    - Replace any non-alphanumeric run with a single underscore
    - Strip leading/trailing underscores
    """
    stem = os.path.splitext(filename)[0]
    key  = stem.lower()
    key  = re.sub(r"[^a-z0-9]+", "_", key)
    key  = key.strip("_")
    return key


def _filename_to_name(filename: str) -> str:
    """
    Convert a filename to a human-readable display name.

    "pn_3_2018_extension_of_time.pdf"
    → "pn_3_2018_extension_of_time"

    Keeps the original stem (without extension) as-is — the user chose that
    filename deliberately. JSON overrides can supply a prettier name.
    """
    return os.path.splitext(filename)[0]


def _discover_folder_sources() -> Dict[str, SourceEntry]:
    """
    Walk knowledge_base/sources/<Category>/<doc_type>/ and return one
    SourceEntry per PDF or CSV file found.

    Folder structure:
        knowledge_base/sources/
            Legislations/
                act/
                    companies_act_2016.pdf
            Practice Notes/
                general/
                    pn_3_2018.pdf
            FAQ/
                faq/
                    faq_incorporation.pdf

    Files in unrecognised category or doc_type folders are skipped with a
    warning so a typo doesn't silently swallow documents.

    Returns {key: SourceEntry}
    """
    discovered: Dict[str, SourceEntry] = {}

    if not os.path.isdir(SOURCES_DIR):
        return discovered

    for cat_name in os.listdir(SOURCES_DIR):
        cat_path = os.path.join(SOURCES_DIR, cat_name)
        if not os.path.isdir(cat_path):
            continue

        if cat_name not in VALID_CATEGORIES:
            logger.warning(
                "  Skipping unknown category folder '%s'. "
                "Valid names: %s", cat_name, VALID_CATEGORIES,
            )
            continue

        for dt_name in os.listdir(cat_path):
            dt_path = os.path.join(cat_path, dt_name)
            if not os.path.isdir(dt_path):
                continue

            if dt_name not in VALID_DOC_TYPES:
                logger.warning(
                    "  Skipping unknown doc_type folder '%s' inside '%s'. "
                    "Valid names: %s", dt_name, cat_name, VALID_DOC_TYPES,
                )
                continue

            _SUPPORTED_EXTENSIONS = (".pdf", ".csv", ".png", ".jpg", ".jpeg", ".webp", ".pptx", ".pptm")

            for fname in sorted(os.listdir(dt_path)):
                if not fname.lower().endswith(_SUPPORTED_EXTENSIONS):
                    continue

                fpath    = os.path.join(dt_path, fname)
                key      = _filename_to_key(fname)
                ext    = os.path.splitext(fname)[1].lower()

                if ext == ".csv":
                    ftype = "csv"
                elif ext in (".pptx", ".pptm"):
                    ftype = "pptx"
                elif ext in (".png", ".jpg", ".jpeg", ".webp"):
                    ftype = "image"
                else:
                    ftype = "pdf"

                if key in discovered:
                    logger.warning(
                        "  Duplicate key '%s' from '%s' — already registered from '%s'. "
                        "Rename one of the files to resolve.",
                        key, fpath, discovered[key].local_path,
                    )
                    continue

                discovered[key] = SourceEntry(
                    key        = key,
                    name       = _filename_to_name(fname),
                    category   = cat_name,
                    type       = ftype,
                    doc_type   = dt_name,
                    enabled    = True,
                    url        = None,
                    local_path = fpath,
                )

    logger.info(
        "Folder discovery: found %d source(s) in '%s'.",
        len(discovered), SOURCES_DIR,
    )
    return discovered


def _load_json_overrides(
    config_path: str = SOURCES_CONFIG,
) -> Dict[str, SourceEntry]:
    """
    Load knowledge_sources.json. Returns {key: SourceEntry}.

    JSON entries can:
    - Add sources not covered by folders (e.g., remote URLs)
    - Override fields of folder-discovered sources (name, enabled, doc_type)
    - Disable a folder-discovered source without deleting its file
    """
    if not os.path.exists(config_path):
        return {}   # JSON is entirely optional

    try:
        with open(config_path, "r", encoding="utf-8") as fh:
            data = json.load(fh)
    except Exception as exc:
        logger.error("Failed to read %s: %s", config_path, exc)
        return {}

    overrides: Dict[str, SourceEntry] = {}
    for raw in data.get("sources", []):
        if "_note" in raw or "_comment" in raw:
            continue
        try:
            key = raw["key"]
            overrides[key] = SourceEntry(
                key        = key,
                name       = raw.get("name", key),
                category   = raw.get("category", ""),
                type       = raw.get("type", "pdf"),
                doc_type   = raw.get("doc_type", "act"),
                enabled    = raw.get("enabled", True),
                url        = raw.get("url"),
                local_path = raw.get("local_path"),
            )
        except KeyError as exc:
            logger.warning(
                "Skipping malformed JSON entry (missing field %s): %s", exc, raw
            )

    logger.info(
        "JSON overrides: loaded %d entry/entries from '%s'.",
        len(overrides), config_path,
    )
    return overrides


def load_sources(config_path: str = SOURCES_CONFIG) -> List[SourceEntry]:
    """
    Discover all sources from two tracks, then merge:

    Track 1 – Folder scan (knowledge_base/sources/<Category>/<doc_type>/)
        Every PDF/CSV file is automatically registered.
        The folder path determines category and doc_type — no config needed.

    Track 2 – JSON overrides (knowledge_sources.json, optional)
        Entries here can:
          • Override any field of a folder-discovered source (name, enabled, …)
          • Add sources with remote URLs not reachable via folder scan
          • Disable a folder source without deleting the file

    Merge rule:
        JSON keys that match a folder key → merged (JSON wins on each field)
        JSON keys with no folder match    → added as standalone sources
        Folder keys with no JSON match    → used as-is (fully automatic)
    """
    folder_sources   = _discover_folder_sources()
    json_overrides   = _load_json_overrides(config_path)

    merged: Dict[str, SourceEntry] = dict(folder_sources)   # start with all folder sources

    for key, json_entry in json_overrides.items():
        if key in merged:
            # Override only fields explicitly set in JSON
            base = merged[key]
            merged[key] = SourceEntry(
                key             = key,
                name            = json_entry.name       if json_entry.name       != key       else base.name,
                category        = json_entry.category   if json_entry.category               else base.category,
                type            = json_entry.type,
                doc_type        = json_entry.doc_type,
                enabled         = json_entry.enabled,
                relates_to_acts = json_entry.relates_to_acts if json_entry.relates_to_acts else base.relates_to_acts,
                url             = json_entry.url        or base.url,
                local_path      = json_entry.local_path or base.local_path,
            )
            logger.debug("  JSON override applied to folder source '%s'.", key)
        else:
            # Standalone JSON source (e.g., remote URL)
            merged[key] = json_entry

    result = [s for s in merged.values() if s.enabled]
    result.sort(key=lambda s: (s.category, s.name))

    total   = len(merged)
    enabled = len(result)
    logger.info(
        "Sources ready: %d total, %d enabled (%d disabled by JSON).",
        total, enabled, total - enabled,
    )
    return result


# ─── Text cleaning ────────────────────────────────────────────────────────────

def clean_text(text: str) -> str:
    """
    Normalise raw PDF text for parsing.

    Approach from the preprocessing notebook, extended with:
    - Unicode dash/quote normalisation (preserves meaning)
    - PDF hyphenation artefact removal (words split across lines)
    - Whitespace normalisation

    Deliberately NOT done here:
    - Lemmatization  (would distort legal terms and section numbers)
    - Non-ASCII stripping  (would remove valid legal symbols like §, ©)
    """
    # ── Dashes ──────────────────────────────────────────────────────────────
    text = text.replace("\u2013", "-")   # en-dash
    text = text.replace("\u2014", "-")   # em-dash
    text = text.replace("\u2012", "-")   # figure-dash

    # ── Quotes ──────────────────────────────────────────────────────────────
    text = text.replace("\u201c", '"').replace("\u201d", '"')   # " "
    text = text.replace("\u2018", "'").replace("\u2019", "'")   # ' '
    text = text.replace("\u201a", "'").replace("\u201b", "'")

    # ── spaCy / NLP pipeline artefacts (if text was pre-processed externally)
    text = text.replace("-LRB-", "(").replace("-RRB-", ")")
    text = text.replace("-LSB-", "[").replace("-RSB-", "]")

    # ── Fix hyphenation split across line breaks: "incor-\nporated" → "incorporated"
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)

    # ── Normalize line endings ───────────────────────────────────────────────
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    # ── Collapse multiple spaces (but keep newlines) ─────────────────────────
    text = re.sub(r"[ \t]{2,}", " ", text)

    # ── Collapse 3+ consecutive blank lines to 2 ────────────────────────────
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()

# ─── OCR via Ollama Vision ────────────────────────────────────────────────────

def ocr_image_bytes(image_bytes: bytes, mime_type: str = "image/png") -> str:
    """
    Send one image to the configured Ollama vision model and return the
    extracted text. Returns "" on any failure.

    Supported models: glm-ocr, llava, minicpm-v, qwen2-vl, etc.
    Set CHATSSM_OCR_MODEL env var to switch models.
    """
    import base64
    b64 = base64.b64encode(image_bytes).decode("ascii")
    try:
        resp = requests.post(
            "http://localhost:11434/api/generate",
            json={
                "model":  OCR_MODEL,
                "prompt": (
                    "Extract ALL text from this image exactly as it appears. "
                    "If the page has multiple columns, read each column top-to-bottom "
                    "before moving to the next column (left column first, then right). "
                    "Preserve paragraph breaks and numbered list structure. "
                    "Output only the extracted text — no commentary, no labels."
                ),
                "images":  [b64],
                "stream":  False,
                "options": {
                    "temperature": 0.0,
                    "num_gpu":     OCR_NUM_GPU,
                }
            },
            timeout=OCR_TIMEOUT,
        )
        if resp.status_code != 200:
            logger.warning("  OCR model returned HTTP %s.", resp.status_code)
            return ""
        text = resp.json().get("response", "").strip()
        # Strip any accidental <think> blocks from reasoning models
        text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL).strip()
        return text
    except requests.Timeout:
        logger.warning("  OCR timed out after %ds for image (%d bytes).", OCR_TIMEOUT, len(image_bytes))
        return ""
    except Exception as exc:
        logger.warning("  OCR failed: %s", exc)
        return ""

# ─── PDF acquisition ──────────────────────────────────────────────────────────

def get_pdf_bytes(source: SourceEntry) -> Optional[bytes]:
    """Download PDF from URL or read from local path."""
    if source.local_path and os.path.exists(source.local_path):
        logger.info("  Reading local PDF: %s", source.local_path)
        try:
            with open(source.local_path, "rb") as fh:
                return fh.read()
        except Exception as exc:
            logger.error("  Failed to read local file: %s", exc)
            return None

    if source.url:
        logger.info("  Downloading PDF from URL …")
        try:
            resp = requests.get(source.url, timeout=60)
            resp.raise_for_status()
            logger.info("  Downloaded %d bytes.", len(resp.content))
            return resp.content
        except requests.Timeout:
            logger.error("  Download timed out.")
        except Exception as exc:
            logger.error("  Download failed: %s", exc)
        return None

    logger.error("  Source '%s' has no url or local_path.", source.key)
    return None


# ─── PDF text extraction ──────────────────────────────────────────────────────

def extract_pdf_text(pdf_bytes: bytes) -> str:
    """
    Extract text from a PDF.

    Per-page strategy:
      1. pdfplumber native extraction.
      2. If a page yields < MIN_TEXT_PER_PAGE chars (scanned/image page),
         render it to PNG via pdf2image and send it to the OCR model.
      3. Tables are always extracted natively as pipe-delimited rows.
      4. If pdfplumber fails entirely, fall back to pypdf (text-only, no OCR).

    Requires pdf2image + poppler for the OCR path. If pdf2image is not
    installed, scanned pages are skipped with a warning.
    """
    try:
        import pdf2image as _pdf2image
        _have_pdf2image = True
    except (ImportError, Exception) as _e:
        _pdf2image = None
        _have_pdf2image = False
        logger.warning(
            "  pdf2image unavailable (%s) — scanned pages will use pypdf fallback. "
            "Fix: pip install pdf2image  and install poppler "
            "(Ubuntu: sudo apt install poppler-utils | "
            "macOS: brew install poppler | "
            "Windows: download from https://github.com/oschwartz10612/poppler-windows)",
            type(_e).__name__,
        )

    pages_out: List[str] = []
    ocr_count  = 0
    total_pages = 0

    # ── pdfplumber (primary) ────────────────────────────────────────────────
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            total_pages = len(pdf.pages)
            native_texts: List[str] = []
            table_texts:  List[str] = []
            needs_ocr:    List[int] = []

            for page_num, page in enumerate(pdf.pages):
                # ── Native text extraction ────────────────────────────────────
                found_tables = page.find_tables()
                table_bboxes = [t.bbox for t in found_tables]

                if table_bboxes:
                    def _outside(obj, bboxes=table_bboxes):
                        x0, top = obj.get("x0", 0), obj.get("top", 0)
                        for (tx0, ttop, tx1, tbottom) in bboxes:
                            if tx0 <= x0 <= tx1 and ttop <= top <= tbottom:
                                return False
                        return True
                    native_text = page.filter(_outside).extract_text() or ""
                else:
                    native_text = page.extract_text() or ""

                # ── Table content ─────────────────────────────────────────────
                table_text = ""
                for table in page.extract_tables() or []:
                    rows = []
                    for row in table:
                        row_text = " | ".join(str(c).strip() if c else "" for c in row)
                        if any(c and str(c).strip() for c in row):
                            rows.append(row_text)
                    if rows:
                        table_text += "\n[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]\n"

                native_texts.append(native_text)
                table_texts.append(table_text)
                if len(native_text) < MIN_TEXT_PER_PAGE:
                    needs_ocr.append(page_num)

        # ── Batch OCR all scanned pages — parallel render + OCR ──────────────────
        ocr_results: Dict[int, str] = {}
        if needs_ocr and _have_pdf2image:
            logger.info(
                "  %d page(s) need OCR. Rendering with %d workers …",
                len(needs_ocr), OCR_WORKERS,
            )

            def _render_and_ocr(page_num: int) -> Tuple[int, str]:
                """Render one page to JPEG and OCR it. Runs in a thread pool."""
                try:
                    from PIL import Image as PILImage
                    page_images = _pdf2image.convert_from_bytes(
                        pdf_bytes,
                        dpi        = 300,
                        fmt        = "jpeg",
                        first_page = page_num + 1,
                        last_page  = page_num + 1,
                    )
                    if not page_images:
                        logger.warning("    Page %d: pdf2image returned no image.", page_num + 1)
                        return page_num, ""
                    img = page_images[0]
                    del page_images

                    max_edge = 1400
                    if max(img.width, img.height) > max_edge:
                        scale = max_edge / max(img.width, img.height)
                        resample_filter = (
                            PILImage.Resampling.LANCZOS
                            if hasattr(PILImage, "Resampling") else 1
                        )
                        img = img.resize(
                            (int(img.width * scale), int(img.height * scale)),
                            resample=resample_filter,
                        )

                    buf = io.BytesIO()
                    img.convert("RGB").save(buf, format="JPEG", quality=80)
                    ocr_text = ocr_image_bytes(buf.getvalue(), "image/jpeg")

                    if not ocr_text:
                        logger.warning(
                            "    Page %d: OCR failed; retrying at half resolution.", page_num + 1
                        )
                        smaller = img.resize((img.width // 2, img.height // 2), resample=1)
                        buf2 = io.BytesIO()
                        smaller.convert("RGB").save(buf2, format="JPEG", quality=75)
                        del smaller
                        ocr_text = ocr_image_bytes(buf2.getvalue(), "image/jpeg")

                    del img
                    return page_num, ocr_text or ""

                except Exception as exc:
                    logger.warning("    Page %d: render or OCR failed: %s", page_num + 1, exc)
                    return page_num, ""

            try:
                with ThreadPoolExecutor(max_workers=OCR_WORKERS) as pool:
                    futures = {pool.submit(_render_and_ocr, pn): pn for pn in needs_ocr}
                    for future in as_completed(futures):
                        pn, text = future.result()
                        if text:
                            ocr_results[pn] = text
                            ocr_count += 1
                        else:
                            logger.warning("    Page %d: OCR returned empty after retry.", pn + 1)
            except Exception as exc:
                logger.warning("  Parallel OCR failed: %s", exc)

        # ── Assemble final pages ──────────────────────────────────────────
        for page_num in range(total_pages):
            native_text = native_texts[page_num]
            table_text  = table_texts[page_num]
            if len(native_text) >= MIN_TEXT_PER_PAGE:
                pages_out.append(native_text + table_text)
            elif page_num in ocr_results:
                # Prepend any real native text (e.g. page headers) before OCR body
                prefix = (native_text.strip() + "\n") if native_text.strip() else ""
                pages_out.append(
                    f"[OCR PAGE {page_num + 1}]\n{prefix}{ocr_results[page_num]}\n{table_text}"
                )
            elif native_text.strip():
                pages_out.append(native_text + table_text)

        full = "\n\n".join(pages_out)
        if ocr_count:
             logger.info(
                "  Extraction complete: %d / %d page(s) via OCR, %d total chars.",
                ocr_count, total_pages, len(full),
            )
             return full
        else:
            logger.info("  pdfplumber: %d chars from %d pages.", len(full), total_pages)

            if len(full) > 500 or "[TABLE]" in full:
                return full
            logger.warning("  pdfplumber returned very little text; trying pypdf.")

    except Exception as exc:
        logger.warning("  pdfplumber failed (%s); trying pypdf.", exc)

    try:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        pages  = [page.extract_text() or "" for page in reader.pages]
        full   = "\n\n".join(p for p in pages if p)
        logger.info("  pypdf fallback: %d chars.", len(full))
        return full
    except Exception as exc:
        logger.error("  pypdf also failed: %s", exc)

    return ""

# ─── Legal document parser ────────────────────────────────────────────────────

def parse_legal_document(
    text: str,
    source_key: str,
    source_name: str,
) -> List[Dict]:
    """
    Parse a cleaned legal document into structured section records.

    Handles the standard Malaysian Act structure:
        PART  ->  Division  ->  Subdivision  ->  Section  ->  Content

    Each returned dict has all CSV_COLUMNS keys.

    Fix log (vs original version)
    ──────────────────────────────
    1. _PART_RE no longer matches "part in the management…"
       (was a false-positive due to re.IGNORECASE on [IVXLCDM])

    2. Marginal note extraction:
       pdfplumber reads the marginal note heading for section N as the last
       line of section N-1. This function now detects that pattern and:
         (a) removes the marginal note from section N-1's content, and
         (b) uses it as the section_title of section N.

    3. Added noise filters for:
       - "Laws of Malaysia ACT 743" running headers
       - "NOTE-The Companies Act 1965 …" editorial footnotes
    """
    rows: List[Dict] = []

    current_part        = "General"
    current_division    = ""
    current_subdivision = ""
    current_section     = ""
    current_title       = ""
    current_content:    List[str] = []
    prev_nonempty       = ""

    # Pending title from a marginal note stripped off the previous section.
    # Set when we detect "Section 23 ends with 'Cessation of partnership interest'"
    # so that Section 24 can use it as its title.
    pending_title       = ""

    def _is_structural(s: str) -> bool:
        return bool(
            _PART_RE.match(s)
            or _DIV_RE.match(s)
            or _SUBDIV_RE.match(s)
            or _SECTION_RE.match(s)
            or _TOC_RE.search(s)
            or _PAGE_RE.match(s)
        )

    def _flush() -> None:
        """
        Save the current section to rows if it has enough content.

        Before saving, check whether the LAST line of current_content is a
        marginal note heading. If so:
          • Remove it from the content (it does not belong there).
          • Store it in pending_title so the NEXT section can use it.
        """
        nonlocal pending_title

        if not current_section:
            return

        content_lines = list(current_content)

        # ── Marginal note detection ──────────────────────────────────────────
        # Look at the last non-empty line of content. If it is a marginal note,
        # pull it off and save it for the next section.
        extracted_note = ""
        for j in range(len(content_lines) - 1, -1, -1):
            last = content_lines[j].strip()
            if not last:
                continue
            if _is_marginal_note(last) and not _is_structural(last):
                extracted_note = last
                content_lines.pop(j)
            break

        content_str = "\n".join(content_lines).strip()
        if len(content_str) < MIN_CONTENT_CHARS:
            # Section too short to save — but preserve the marginal note
            if extracted_note:
                pending_title = extracted_note
            return

        rows.append({
            "source_key":    source_key,
            "source_name":   source_name,
            "part":          current_part,
            "division":      current_division,
            "subdivision":   current_subdivision,
            "section":       current_section,
            "section_title": current_title,
            "content":       content_str,
        })

        # Pass the extracted marginal note to the next section
        if extracted_note:
            pending_title = extracted_note

    lines = text.splitlines()
    i = 0

    while i < len(lines):
        line = lines[i].strip()

        # ── Skip noise ───────────────────────────────────────────────────────
        if not line:
            i += 1
            continue
        if _TOC_RE.search(line):
            i += 1
            continue
        if _PAGE_RE.match(line):
            i += 1
            continue
        if _GAZETTE_RE.match(line):
            i += 1
            continue
        if _LAWS_HDR_RE.match(line):      # "Laws of Malaysia ACT 743"
            i += 1
            continue
        if _NOTE_RE.match(line):          # "NOTE-The Companies Act 1965…"
            i += 1
            continue
        if _TABLE_TAG_RE.match(line):  
            i += 1
            continue

        # ── PART header ──────────────────────────────────────────────────────
        if _PART_RE.match(line):
            _flush()
            current_section     = ""
            current_title       = ""
            current_content     = []
            current_division    = ""
            current_subdivision = ""

            # Grab the Part title from the next non-empty line if it looks
            # like a title (not another structural header or TOC entry).
            part_label = line
            j = i + 1
            while j < len(lines) and not lines[j].strip():
                j += 1
            if j < len(lines):
                peek = lines[j].strip()
                if (peek
                        and not _PART_RE.match(peek)
                        and not _SECTION_RE.match(peek)
                        and not _TOC_RE.search(peek)
                        and not _PAGE_RE.match(peek)
                        and not _LAWS_HDR_RE.match(peek)):
                    part_label = f"{line}: {peek}"
                    i = j

            current_part = part_label
            i += 1
            continue

        # ── Division header ───────────────────────────────────────────────────
        if _DIV_RE.match(line):
            _flush()
            current_section     = ""
            current_title       = ""
            current_content     = []
            current_subdivision = ""

            div_label = line
            j = i + 1
            while j < len(lines) and not lines[j].strip():
                j += 1
            if j < len(lines):
                peek = lines[j].strip()
                if (peek
                        and not _PART_RE.match(peek)
                        and not _DIV_RE.match(peek)
                        and not _SECTION_RE.match(peek)
                        and not _TOC_RE.search(peek)):
                    div_label = f"{line}: {peek}"
                    i = j

            current_division = div_label
            i += 1
            continue

        # ── Subdivision header ────────────────────────────────────────────────
        if _SUBDIV_RE.match(line):
            _flush()
            current_section = ""
            current_title   = ""
            current_content = []

            sdiv_label = line
            j = i + 1
            while j < len(lines) and not lines[j].strip():
                j += 1
            if j < len(lines):
                peek = lines[j].strip()
                if (peek
                        and not _PART_RE.match(peek)
                        and not _DIV_RE.match(peek)
                        and not _SUBDIV_RE.match(peek)
                        and not _SECTION_RE.match(peek)
                        and not _TOC_RE.search(peek)):
                    sdiv_label = f"{line}: {peek}"
                    i = j

            current_subdivision = sdiv_label
            i += 1
            continue

        # ── Section header ────────────────────────────────────────────────────
        sm = _SECTION_RE.match(line)
        if sm:
            _flush()   # this may set pending_title from the marginal note

            section_num  = sm.group(1).strip()
            inline_title = sm.group(2).strip()

            current_section = f"Section {section_num}"
            current_content = []

            # Title priority:
            #   1. Inline title on the same line as the section number
            #   2. Marginal note extracted from the previous section (pending_title)
            #   3. Last non-structural line before the section number (prev_nonempty)
            if inline_title:
                current_title = inline_title
                pending_title = ""
            elif pending_title:
                current_title = pending_title
                pending_title = ""
            elif (prev_nonempty
                    and not _is_structural(prev_nonempty)
                    and not _LAWS_HDR_RE.match(prev_nonempty)):
                current_title = prev_nonempty
            else:
                current_title = ""

            i += 1
            continue

        # ── Regular content ───────────────────────────────────────────────────
        if current_section:
            current_content.append(line)

        prev_nonempty = line
        i += 1

    # Flush the final section
    _flush()

    logger.info("  Parsed %d sections from '%s'.", len(rows), source_name)
    return rows


# ─── General document parser (Practice Notes, Guidelines, Circulars) ──────────
#
# SSM Practice Note structure (from real document analysis):
#
#   TITLE LINE(S)                         ← document title, free text
#   OBJECTIVE / BACKGROUND / etc.         ← ALL-CAPS section headings  ← CHUNK BOUNDARY
#   A. Lettered Sub-heading               ← lettered subsections        ← CHUNK BOUNDARY
#   B. Another Sub-heading
#   2.  Numbered paragraph content.       ← numbered paragraphs         ← CONTENT, not boundary
#   3.  More content.
#   (a) Sub-item content.                 ← lettered sub-items          ← CONTENT
#   (b) More sub-items.
#
# CHUNKING STRATEGY:
#   Only ALL-CAPS headings and lettered subsections (A., B.) are chunk
#   boundaries.  Numbered paragraphs (2., 3., 10.) are accumulated as
#   content WITHIN the current chunk.
#
#   This produces semantically coherent chunks like:
#     Chunk: "BACKGROUND" → paragraphs 3–7 (all context about background)
#     Chunk: "A. Initial Application" → paragraphs 8–9 (the rule itself)
#   Rather than tiny 1-sentence chunks per paragraph number.
#
#   If a chunk grows beyond MAX_GENERAL_CHUNK_CHARS, it is split at
#   sentence boundaries to stay within the embedding model's context.

# ALL-CAPS heading: ≥ 4 chars, may contain spaces, /, &, (, ), digits (years like 2016)
# e.g. "OBJECTIVE", "BACKGROUND", "CIRCUMSTANCES WHERE COMPANY CAN APPLY FOR EOT UNDER THE COMPANIES ACT 2016"
# Must be ENTIRELY uppercase (or with permitted symbols) — excludes mixed-case content
_CAPS_HEADING_RE = re.compile(r"^[A-Z][A-Z\s/&\-\(\)\d]{3,}$")

# Lettered subsection: "A.", "A. Title", "B.", "B. Long Title Here"
# Deliberately requires a SINGLE uppercase letter to avoid matching content like
# "(a) Private Company" — those are sub-items, not headings
_LETTERED_SUBSECTION_RE = re.compile(r"^([A-Z])\.\s+\S")

# Numbered paragraph: "2.", "10.", "2.1", "1.1.1"
# Used to STRIP the number prefix so content is clean, NOT used as chunk boundary
_PARA_NUM_STRIP_RE = re.compile(r"^(\d+(?:\.\d+){0,3})[.)]\s+")

# Sub-item: "(a)", "(b)", "(i)", "(ii)", "(iii)"
_SUB_ITEM_RE = re.compile(r"^\([a-z]+\)\s+")

# Appendix / Annex / Schedule — treated as a major chunk boundary
_APPENDIX_RE = re.compile(r"^(Appendix|Annex|Schedule|Attachment)\s+\S+", re.IGNORECASE)

# Practice Note page header/footer noise:
# "1 PN 3/2018 26 July 2018", "2 PN 3/2018 26 July 2018"
# Pattern: optional digits, then "PN" + number/year + date
_PN_HEADER_RE = re.compile(r"^\d*\s*PN\s+\d+/\d{4}", re.IGNORECASE)

# Maximum characters per chunk before forcing a sentence-boundary split
MAX_GENERAL_CHUNK_CHARS = 1500


def _is_major_boundary(line: str) -> bool:
    """
    Return True only for lines that should FLUSH the current chunk and start a new one.

    For SSM Practice Notes this means:
      - ALL-CAPS section headings  (OBJECTIVE, BACKGROUND, REQUIREMENTS, ...)
      - Lettered subsections       (A. Initial Application, B. EOT Application, ...)
      - Appendix / Annex / Schedule lines
      - Title-case short headings that are NOT numbered paragraphs

    Numbered paragraphs (2., 3., 10.) are NOT chunk boundaries — they are
    accumulated as content within the current chunk.
    """
    if not line or len(line) < 2:
        return False
    if line.endswith(",") or line.endswith(";"):
        return False
    if _CAPS_HEADING_RE.match(line) and not _PAGE_RE.match(line):
        return True
    if _LETTERED_SUBSECTION_RE.match(line):
        return True
    if _APPENDIX_RE.match(line):
        return True
    if _PARA_LABEL_RE.match(line):   # "Para 3", "Paragraph 5"
        return True
    return False


# Keep _PARA_LABEL_RE for Appendix/Para matching
_PARA_LABEL_RE = re.compile(
    r"^(Para(?:graph)?\s+\d+|Appendix\s+\S+|Annex\s+\S+)\b", re.IGNORECASE
)


def parse_general_document(
    text: str,
    source_key: str,
    source_name: str,
) -> List[Dict]:
    """
    Heading-aware chunker for SSM Practice Notes, Guidelines, and Circulars.

    Chunk boundaries
    ----------------
    Only ALL-CAPS headings and lettered subsections (A., B.) start new chunks.
    Numbered paragraphs (2., 3., 10.) are CONTENT within the current chunk —
    they are included with the heading that precedes them.

    Example output for the practice note format
    -------------------------------------------
    Row 1:
      section  = "OBJECTIVE"
      content  = "2. This Practice Note serves to clarify on the application
                  for extension of time under subsections 609(2) ..."

    Row 2:
      section  = "BACKGROUND"
      content  = "3. There is a need to lodge documents in a timely manner ...
                  [paragraphs 3–7 all grouped here]"

    Row 3:
      section  = "A. Initial Application and Subsequent Approvals for EOT"
      content  = "8. Application must be received by the Registrar at least
                  seven (7) days ... [paragraphs 8–9]"

    Oversized chunks (> MAX_GENERAL_CHUNK_CHARS) are split at sentence
    boundaries, with the section heading prepended to each sub-chunk so
    every chunk is self-contained for retrieval.
    """
    rows: List[Dict] = []

    current_heading     = ""     # current chunk's heading (the boundary line)
    current_part        = ""     # most recent ALL-CAPS heading for the 'part' column
    current_lines: List[str] = []

    def _flush() -> None:
        nonlocal current_lines
        content = "\n".join(current_lines).strip()
        current_lines = []
        if len(content) < MIN_CONTENT_CHARS:
            return

        if len(content) > MAX_GENERAL_CHUNK_CHARS:
            sub_chunks = _split_at_sentences(content, MAX_GENERAL_CHUNK_CHARS)
        else:
            sub_chunks = [content]

        for chunk in sub_chunks:
            chunk = chunk.strip()
            if len(chunk) < MIN_CONTENT_CHARS:
                continue
            rows.append({
                "source_key":    source_key,
                "source_name":   source_name,
                "part":          current_part,
                "division":      "",
                "subdivision":   "",
                "section":       current_heading,
                "section_title": current_heading,
                "content":       chunk,
            })

    for raw_line in text.splitlines():
        line = raw_line.strip()

        # ── Noise filters ─────────────────────────────────────────────────────
        if not line:
            continue
        if _TOC_RE.search(line):
            continue
        if _PAGE_RE.match(line):
            continue
        if _GAZETTE_RE.match(line):
            continue
        if _LAWS_HDR_RE.match(line):
            continue
        if _NOTE_RE.match(line):
            continue
        if _PN_HEADER_RE.match(line):       # "1 PN 3/2018 26 July 2018"
            continue
        if _TABLE_TAG_RE.match(line):
            continue

        # ── Major chunk boundary ──────────────────────────────────────────────
        if _is_major_boundary(line):
            _flush()
            current_heading = line
            # Promote ALL-CAPS headings to the 'part' column
            if _CAPS_HEADING_RE.match(line):
                current_part = line
            continue

        # ── Numbered paragraph: strip the number, keep the text ───────────────
        # "2.  This Practice Note serves to..." → "This Practice Note serves to..."
        # The number is included in the content string as-is for citation purposes
        # We keep the original line (with number) so the LLM can cite "paragraph 2"
        current_lines.append(line)

    # Flush the final chunk
    _flush()

    logger.info("  Parsed %d chunks from '%s'.", len(rows), source_name)
    return rows


def _split_at_sentences(text: str, max_chars: int) -> List[str]:
    """
    Split *text* into chunks ≤ max_chars at sentence boundaries.
    Falls back to hard split if no sentence boundary is found.
    """
    if len(text) <= max_chars:
        return [text]
    
    if "[TABLE]" in text or text.count("|") > 3:
        return [text] 

    chunks: List[str] = []
    while len(text) > max_chars:
        # Find the last sentence-ending punctuation before the limit
        split_at = -1
        for m in re.finditer(r"[.!?]\s+", text[:max_chars]):
            candidate = m.end()
            before = text[:candidate]
            if "|" not in before or (
                "\n" in before and before.rindex("\n") > before.rindex("|")
            ):
                split_at = candidate
        if split_at == -1:
            # No sentence boundary — split at last whitespace
            split_at = text[:max_chars].rfind(" ")
        if split_at <= 0:
            split_at = max_chars   # hard cut

        chunks.append(text[:split_at].strip())
        text = text[split_at:].strip()

    if text:
        chunks.append(text)
    return chunks


# ─── FAQ document parser ───────────────────────────────────────────────────────

# "Q:" / "A:" style
_FAQ_Q_COLON   = re.compile(r"^Q\s*[:\.]\s*(.+)", re.IGNORECASE)
_FAQ_A_COLON   = re.compile(r"^A\s*[:\.]\s*(.+)", re.IGNORECASE)

# Numbered question: "1.", "1)", "(1)", "Q1.", "Q1:"
_FAQ_Q_NUM     = re.compile(
    r"^(?:Q\.?\s*)?(\d+)[.)]\s+(.+)", re.IGNORECASE
)

# "Question 1:", "Question:" (unnumbered)
_FAQ_Q_WORD    = re.compile(r"^Question\s*(?:\d+)?\s*[:.]?\s*(.+)?", re.IGNORECASE)
_FAQ_A_WORD    = re.compile(r"^Answer\s*[:.]?\s*(.+)?",               re.IGNORECASE)

# FAQ page footer noise:
# "PART A_amended 100323 r2 – Page 1"  "PART B_amended – Page 2"
# Pattern: PART <letter>_ followed by anything
_FAQ_PAGE_FOOTER_RE = re.compile(r"^PART\s+[A-Z]_", re.IGNORECASE)

# Table tag lines produced by extract_pdf_text()
_TABLE_TAG_RE = re.compile(r"^\[(TABLE|/TABLE)\]$", re.IGNORECASE)

# Title-case section heading inside FAQ body
# e.g. "Transitional Provisions Relating to Abolition of Nominal Value"
# Heuristic: starts Capital, contains mostly title-case words, no terminal punctuation,
# ≤ 15 words, no digits at the start
_FAQ_TITLE_HEADING_RE = re.compile(
    r"^[A-Z][a-zA-Z]+((\s+[A-Za-z]+){1,14})$"
)


def _is_faq_section_heading(line: str) -> bool:
    """
    Return True if *line* looks like a section heading within a FAQ document.
    Catches both ALL-CAPS headings and Title Case multi-word headings that
    appear between Q/A groups (e.g. "Transitional Provisions Relating to...").
    """
    if not line or len(line) < 4:
        return False
    if line.endswith("?") or line.endswith(".") or line.endswith(":"):
        return False
    # ALL-CAPS headings already caught by _is_major_boundary — included here
    # too so this function is self-contained for use in the FAQ parser
    if _CAPS_HEADING_RE.match(line) and not _PAGE_RE.match(line):
        return True
    # Title Case heading: no digits in it, not too short, not a question
    if (_FAQ_TITLE_HEADING_RE.match(line)
            and not any(c.isdigit() for c in line)
            and len(line.split()) >= 3):
        return True
    return False


def parse_faq_document(
    text: str,
    source_key: str,
    source_name: str,
) -> List[Dict]:
    """
    FAQ-specific parser that detects Q/A pairs and keeps each pair together
    as one self-contained chunk.

    Supported FAQ formats
    ---------------------
    Style 1 – Colon-prefix:        Q: <question>    A: <answer>
    Style 2 – Numbered questions:  1. <question>    Answer: <answer>
    Style 3 – Word prefix:         Question: ...    Answer: ...
    Style 4 – Plain question:      <sentence ending in ?>
                                   Answer: <answer>
               (used in this SSM FAQ — no Q: or number prefix)

    Each Q+A pair is stored as:
        "Q: <question text>\n\nA: <answer text>"

    Table handling
    --------------
    pdfplumber wraps tables in [TABLE]...[/TABLE] tags.
    The [TABLE]/[/TABLE] wrapper lines are stripped.
    The pipe-delimited row content is kept — it contains useful information
    (e.g. terminology tables, lodgement timelines) that the LLM needs.

    Page footer noise
    -----------------
    "PART A_amended 100323 r2 – Page 1" lines are filtered out.

    Section headings within FAQ
    ---------------------------
    Title-case headings between Q/A groups (e.g. "Transitional Provisions
    Relating to Abolition of Nominal Value") are detected and stored in the
    'part' column so the LLM knows which topic area each Q/A belongs to.
    """
    rows: List[Dict] = []

    current_section  = ""
    current_q        = ""
    current_a_lines: List[str] = []
    mode             = "scan"    # "scan" | "in_q" | "in_a"
    prev_line        = ""        # used to detect plain-question format

    def _flush_qa() -> None:
        nonlocal current_q, current_a_lines, mode
        if not current_q:
            return
        answer = "\n".join(current_a_lines).strip()
        if not answer:
            if len(current_q) >= MIN_CONTENT_CHARS:
                rows.append({
                    "source_key":    source_key,
                    "source_name":   source_name,
                    "part":          current_section,
                    "division":      "",
                    "subdivision":   "",
                    "section":       current_section,
                    "section_title": current_q[:120],
                    "content":       current_q,
                })
        else:
            content = f"Q: {current_q}\n\nA: {answer}"
            rows.append({
                "source_key":    source_key,
                "source_name":   source_name,
                "part":          current_section,
                "division":      "",
                "subdivision":   "",
                "section":       current_section,
                "section_title": current_q[:120],
                "content":       content,
            })
        current_q       = ""
        current_a_lines = []
        mode            = "scan"

    for raw_line in text.splitlines():
        line = raw_line.strip()

        # ── Noise filters ─────────────────────────────────────────────────────
        if not line:
            continue
        if _TOC_RE.search(line):
            continue
        if _PAGE_RE.match(line):
            continue
        if _GAZETTE_RE.match(line):
            continue
        if _LAWS_HDR_RE.match(line):
            continue
        if _NOTE_RE.match(line):
            continue
        if _PN_HEADER_RE.match(line):
            continue
        if _FAQ_PAGE_FOOTER_RE.match(line):       # "PART A_amended 100323 r2 – Page 1"
            continue
        if _TABLE_TAG_RE.match(line):             # "[TABLE]" / "[/TABLE]" wrapper tags
            continue                              # keep row content, strip tags

        # ── Section / heading ─────────────────────────────────────────────────
        if _is_faq_section_heading(line) and not _FAQ_Q_NUM.match(line):
            _flush_qa()
            current_section = line
            mode = "scan"
            prev_line = line
            continue

        # ── Q: / Question: / Numbered question ────────────────────────────────
        m_qc = _FAQ_Q_COLON.match(line)
        m_qw = _FAQ_Q_WORD.match(line)
        m_qn = _FAQ_Q_NUM.match(line)

        if m_qc or m_qw or m_qn:
            _flush_qa()
            if m_qc:   current_q = m_qc.group(1).strip()
            elif m_qw: current_q = (m_qw.group(1) or "").strip()
            elif m_qn: current_q = m_qn.group(2).strip()
            mode = "in_q"
            prev_line = line
            continue

        # ── A: / Answer: ──────────────────────────────────────────────────────
        m_ac = _FAQ_A_COLON.match(line)
        m_aw = _FAQ_A_WORD.match(line)

        if m_ac or m_aw:
            # ── Style 4: plain question (no prefix) ───────────────────────────
            # If we're in scan mode and the previous non-empty line was a
            # question (ends with ?), treat it as the question for this answer.
            # Example:
            #   "Please clarify if the entire Act will take effect on 31 January?"
            #   "Answer: Once enforced, all provisions will take effect..."
            if mode == "scan" and prev_line.endswith("?"):
                _flush_qa()
                current_q = prev_line

            if m_ac:
                first_line = m_ac.group(1).strip()
            else:
                first_line = (m_aw.group(1) or "").strip() if m_aw else ""
            if first_line:
                current_a_lines.append(first_line)
            mode = "in_a"
            prev_line = line
            continue

        # ── Continuation lines ────────────────────────────────────────────────
        if mode == "in_q":
            current_q = (current_q + " " + line).strip()
        elif mode == "in_a":
            current_a_lines.append(line)
        else:
            # Scan mode — plain content line (narrative paragraph, not Q/A)
            if len(line) >= MIN_CONTENT_CHARS:
                rows.append({
                    "source_key":    source_key,
                    "source_name":   source_name,
                    "part":          current_section,
                    "division":      "",
                    "subdivision":   "",
                    "section":       current_section,
                    "section_title": "",
                    "content":       line,
                })

        prev_line = line

    # Flush the last Q/A pair
    _flush_qa()

    logger.info("  Parsed %d Q/A chunks from '%s'.", len(rows), source_name)
    return rows


# ─── Gazette / Subsidiary Legislation parser ───────────────────────────────
#
# Handles: Companies Gazette, LLP Gazette, Business Registration Gazette,
# and any Malaysian subsidiary legislation (P.U.(A) gazette orders).
#
# Structure:
#   TITLE (bilingual header)
#   ARRANGEMENT OF GAZETTES  (bilingual TOC — must be stripped)
#   Title-case heading           ← chunk boundary ("Citation and commencement")
#   N. (1) Content...            ← gazette body (numbered paragraphs = CONTENT)
#   SCHEDULE                     ← fee table — one chunk per schedule
#   [TABLE] rows [/TABLE]        ← pdfplumber table output
#
# Key differences from Acts:
#   1. "Section" → "Gazette" numbering ("3." not "Section 3.")
#   2. Chunk boundaries are Title Case headings, not PART/Section markers
#   3. Document is BILINGUAL (Malay then English). We keep English only.
#   4. Running header "P.U. (A) 37" appears on every page — must be filtered.
#   5. Bilingual TOC lines must be filtered.
#   6. Fee schedule table rows (pipe-delimited) must be kept together.

# Running header: "P.U. (A) 37", "P.U. (A) 173/1966"
_PU_HDR_RE = re.compile(r"^P\.U\.\s*\([AB]\)\s*[\d/]+", re.IGNORECASE)

# Bilingual TOC line: "Peraturan 1. Nama..." or "Regulation 1. Citation..."
_REG_TOC_RE = re.compile(r"^(Peraturan|Regulation)\s+\d+\.", re.IGNORECASE)

# Bilingual gazette header lines
_GAZETTE_HEADER_RE = re.compile(
    r"^(WARTA KERAJAAN|FEDERAL GOVERNMENT GAZETTE|PERATURAN-PERATURAN|"
    r"SUSUNAN PERATURAN|ARRANGEMENT OF REGULATIONS|"
    r"AKTA SYARIKAT|COMPANIES ACT \d{4}|COMPANIES REGULATIONS \d{4}|"
    r"FI$|FEES$|"                        # ALL-CAPS only — "Fees" (mixed-case) is a valid heading
    r"\(1\) Bil\.|\(1\) Item|"
    r"\(2\) Perkara|\(2\) Matter|"
    r"\(3\) Fi|\(3\) Fee|"
    r"DISIARKAN|PUBLISHED BY|JABATAN PEGUAM|ATTORNEY GENERAL)",
)                                        # deliberately NO re.IGNORECASE — "FI$" must be exact

# Malay-only line detector: contains common Malay legal words with NO English equivalent
# on the same line. This strips the Malay half of each bilingual pair.
_MALAY_LINE_RE = re.compile(
    r"\b(hendaklah|boleh|adalah|dengan|kepada|daripada|atau|bagi|oleh|"
    r"yang|tidak|dan|dalam|ini|itu|pada|telah|akan|jika|mana-mana|"
    r"penyerahsimpanan|pemfailan|seksyen|akta|syarikat|peraturan|"
    r"menteri|pendaftar|rayuan|permohonan|pembatalan|ertinya|bolehlah|"
    r"berkuat|kuasa|dinamakan|diberikan|diarahkan|diserah|simpan|"
    r"dikehendaki|dinyatakan|terkilan|terbatal|dibayar|ditarik)\b",
    re.IGNORECASE,
)
_ENGLISH_ANCHOR_RE = re.compile(
    r"\b(the|shall|must|may|any|all|under|section|regulation|act|company|"
    r"minister|registrar|document|application|appeal|electronic|filing)\b",
    re.IGNORECASE,
)

# Regulation number line: "3.  All documents..." or "8. (1) The fees..."
_GAZETTE_BODY_RE = re.compile(r"^(\d+[A-Z]?)\.\s+(.+)")

# Schedule/Appendix heading
_SCHEDULE_RE = re.compile(r"^(SCHEDULE|JADUAL|Appendix|Annex)\b", re.IGNORECASE)

# Fee table row: contains RM amounts or numeric fee data
_FEE_ROW_RE = re.compile(r"(\d{1,3}(,\d{3})*\.\d{2}|Tiada\b)", re.IGNORECASE)

# Title-case gazette heading — max 90 chars to cover long Malaysian headings
# e.g. "Lodgement of documents or application through electronic filing system" (70 chars)
_GAZETTE_HEADING_RE = re.compile(
    r"^[A-Z][a-zA-Z][\w\s,\-/()&']{1,78}$"
)

# Prose-starter words that never open a gazette heading
_GAZETTE_PROSE_STARTER_RE = re.compile(
    r"^(Pursuant|Subject|Notwithstanding|In\s+accordance|Provided|However|"
    r"Upon|All\s+persons?|Any\s+persons?|Every\s+persons?|No\s+persons?|"
    r"The\s+(Minister|Registrar|company|licensee|holder|applicant|person))\b",
    re.IGNORECASE,
)

_ENGLISH_START   = re.compile(r"^IN\s+exercise\s+of",                    re.IGNORECASE)
_MALAY_START     = re.compile(r"^PADA\s+menjalankan",                     re.IGNORECASE)
_GAZETTE_SIGNING = re.compile(r"^(Dibuat|Made\s+\d|DATO|Yang\s+Berhormat)", re.IGNORECASE)
_SOLO_FEE        = re.compile(r"^[\d,]+\.\d{2}\b")
_SCHED_COL_HDR   = re.compile(
    r"^\(Regulation \d+\)$"
    r"|^FEES?$"
    r"|\(\d+\)\s*(Item|Matter|Fee|Penalty|Bil|Perkara|Remarks|No\.?|Description|Amount)"
    r"|^(Item\s*No\.?|Description(\s+of\s+Matter)?|Amount\s*\(?RM\)?|Fee\s*\(?RM\)?)$",
    re.IGNORECASE,
)

def _is_gazette_heading(line: str) -> bool:
    """
    Return True if line is a gazette section heading — the Title Case
    heading that precedes each numbered gazette body.

    Examples:
      "Citation and commencement"       → True
      "Lodgement of documents..."       → True  (short enough)
      "Fees"                            → True
      "3.  All documents required..."   → False  (gazette body)
      "(1) Upon receipt..."             → False  (sub-gazette body)
    """
    if not line or len(line) < 3:
        return False
    # Must start with a capital letter followed by lowercase
    if not re.match(r"^[A-Z][a-z]", line):
        return False
    # Must not look like a gazette body line
    if _GAZETTE_BODY_RE.match(line):
        return False
    # Must not end with sentence punctuation (headings don't)
    if line.endswith(".") and len(line) > 40:
        return False
    # Must not be a sub-item
    if re.match(r"^\([a-z]\)", line):
        return False
    # Reject long prose sentences
    if len(line) > 100:
        return False
    if len(line.split()) > 12:          # real headings are concise noun phrases
        return False
    if _GAZETTE_PROSE_STARTER_RE.match(line):   # common sentence openers
        return False
    if line.endswith(",") or line.endswith(";"):
        return False
    if _GAZETTE_HEADING_RE.match(line):
        return True
    return False


def _is_malay_line(line: str) -> bool:
    """
    Return True if a line is predominantly Malay with no significant English content.
    Used to strip the Malay half of bilingual gazette documents.

    Strategy: count Malay keyword hits vs English keyword hits.
    If Malay hits > English hits AND no strong English anchors → Malay line.
    """
    malay_hits   = len(_MALAY_LINE_RE.findall(line))
    english_hits = len(_ENGLISH_ANCHOR_RE.findall(line))

    if malay_hits == 0:
        return False                    # no Malay keywords → keep
    if english_hits >= malay_hits:
        return False                    # equal or more English → keep
    if english_hits >= 3:
        return False                    # enough English anchors → keep (bilingual line)
    return True                         # mostly Malay → strip


def parse_gazette_document(
    text: str,
    source_key: str,
    source_name: str,
) -> List[Dict]:
    """
    Parser for Malaysian subsidiary legislation:
    Companies Gazettes, LLP Gazettes, Business Registration Gazettes,
    gazette orders, and similar P.U.(A) documents.

    Chunking strategy
    -----------------
    Each Title Case heading (e.g. "Citation and commencement", "Fees") starts
    a new chunk. All numbered gazette paragraphs under that heading are
    accumulated as content within that chunk.

    The fee SCHEDULE is treated as one dedicated chunk per schedule section,
    with all table rows (pipe-delimited) kept together so the LLM can answer
    fee queries accurately.

    Bilingual handling
    ------------------
    Malaysian gazette documents contain the full Malay text followed by the
    full English text (or interleaved). This parser keeps English lines and
    discards predominantly-Malay lines so chunks don't contain duplicate
    content in two languages.

    Table handling
    --------------
    [TABLE]/[/TABLE] wrapper tags from pdfplumber are stripped.
    Pipe-delimited row content is kept — it contains fee amounts, section
    references, and application types that are essential for queries like
    "what is the fee to incorporate a company?"
    """
    rows: List[Dict] = []

    current_heading     = ""
    current_part        = ""    # running header for part column
    current_lines: List[str] = []
    in_schedule         = False
    schedule_name       = ""

    def _flush() -> None:
        nonlocal current_lines
        content = "\n".join(l for l in current_lines if l).strip()
        current_lines = []
        if len(content) < MIN_CONTENT_CHARS:
            return

        if len(content) > MAX_GENERAL_CHUNK_CHARS:
            sub_chunks = _split_at_sentences(content, MAX_GENERAL_CHUNK_CHARS)
        else:
            sub_chunks = [content]

        for chunk in sub_chunks:
            chunk = chunk.strip()
            if len(chunk) < MIN_CONTENT_CHARS:
                continue
            rows.append({
                "source_key":    source_key,
                "source_name":   source_name,
                "part":          current_part or schedule_name,
                "division":      "",
                "subdivision":   "",
                "section":       current_heading,
                "section_title": current_heading,
                "content":       chunk,
            })

    in_english_body = False   # True once "IN exercise of the powers..." seen
    in_malay_body   = False   # True once "PADA menjalankan..." seen

    for raw_line in text.splitlines():
        line = raw_line.strip()

        # ── Bilingual body boundary detection ─────────────────────────────────
        if _MALAY_START.match(line):
            in_malay_body = True; in_english_body = False
            continue
        if _ENGLISH_START.match(line):
            _flush()
            in_english_body = True; in_malay_body = False
            continue

        # ── Skip the entire Malay body (first half of document) ───────────────
        # The English body is a complete translation — no need for Malay content.
        if in_malay_body and not in_english_body:
            continue

        # ── Universal noise filters ────────────────────────────────────────────
        if not line:
            continue
        if _TOC_RE.search(line):
            continue
        if _PAGE_RE.match(line):
            continue
        if _PU_HDR_RE.match(line):                  # "P.U. (A) 37" running header
            continue
        if _GAZETTE_HEADER_RE.match(line):           # gazette/header boilerplate
            continue
        if _REG_TOC_RE.match(line):                  # "Regulation 1. Citation..."
            continue
        if _TABLE_TAG_RE.match(line):                # "[TABLE]" / "[/TABLE]"
            continue
        if _LAWS_HDR_RE.match(line):
            continue
        if _NOTE_RE.match(line):
            continue
        if _GAZETTE_SIGNING.match(line):                     # "Dibuat ...", "DATO' ..."
            continue

        # ── Bilingual line filter (catches any residual Malay lines) ──────────
        if _is_malay_line(line):
            continue

        # ── Schedule / fee table section ───────────────────────────────────────
        if _SCHEDULE_RE.match(line) and in_english_body:
            _flush()
            in_schedule   = True
            schedule_label  = line.strip()
            schedule_name   = schedule_label
            current_heading = schedule_label
            current_part    = schedule_label
            continue

        # ── Within-schedule sub-heading (e.g. "FEES", "PENALTIES") ───────────
        if in_schedule and _CAPS_HEADING_RE.match(line) and not _SCHED_COL_HDR.match(line):
            current_heading = f"{schedule_name} - {line.strip()}"
            continue

        # ── Fee amount-only lines: attach to preceding item ───────────────────
        # pdfplumber often puts the fee on its own line after the item description.
        # Keep it as content so the LLM knows the actual amount.
        # (Nothing to skip — just let it fall through to current_lines.append)

        # ── Gazette heading (Title Case, not numbered) ─────────────────────
        if _is_gazette_heading(line) and not in_schedule:
            _flush()
            current_heading    = line
            if len(line.split()) <= 8:
                current_part = line
            continue

        # ── Content line ───────────────────────────────────────────────────────
        current_lines.append(line)

    _flush()

    logger.info("  Parsed %d chunks from gazette '%s'.", len(rows), source_name)
    return rows

# ─── PPTX extraction ──────────────────────────────────────────────────────────

def extract_pptx_content(pptx_bytes: bytes) -> List[Dict]:
    """
    Extract structured content from a PowerPoint file.

    Per slide, collects:
      - title text (from TITLE/CENTER_TITLE placeholder)
      - body text (from all other text frames, in reading order)
      - table content as pipe-delimited rows
      - speaker notes
      - raw image bytes (for OCR when a slide has no extractable text)

    Returns a list of slide dicts — one per slide. Empty slides
    (no text, no images) are included so slide numbering is preserved;
    parse_slide_document() will filter them.
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("  python-pptx not installed. Run:  pip install python-pptx")
        return []

    # Resolve shape-type enums — name changed in python-pptx 1.x (PP_PLACEHOLDER_TYPE)
    # Fall back to raw integers if neither name is importable.
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as _PPH
        _TITLE_TYPES = (_PPH.TITLE, _PPH.CENTER_TITLE)
    except ImportError:
        try:
            from pptx.enum.shapes import PP_PLACEHOLDER as _PPH   # 0.x name
            _TITLE_TYPES = (_PPH.TITLE, _PPH.CENTER_TITLE)
        except ImportError:
            _TITLE_TYPES = (1, 3)   # TITLE=1, CENTER_TITLE=3 (stable integer values)

    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
        _PICTURE_TYPE = _MSO.PICTURE
    except ImportError:
        _PICTURE_TYPE = 13          # MSO_SHAPE_TYPE.PICTURE=13 (stable integer value)

    slides_data: List[Dict] = []
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        for slide_num, slide in enumerate(prs.slides, start=1):
            title       = ""
            body_parts: List[str] = []
            table_rows: List[str] = []
            image_blobs: List[bytes] = []

            for shape in slide.shapes:

                # ── Text frames ───────────────────────────────────────────────
                if shape.has_text_frame:
                    lines = []
                    for para in shape.text_frame.paragraphs:
                        para_text = "".join(run.text for run in para.runs).strip()
                        if para_text:
                            lines.append(para_text)
                    combined = "\n".join(lines).strip()
                    if not combined:
                        continue

                    # Is this shape a title placeholder?
                    is_title = False
                    try:
                        if (shape.is_placeholder
                                and shape.placeholder_format.type
                                    in _TITLE_TYPES):
                            is_title = True
                    except Exception:
                        pass

                    if is_title:
                        title = combined
                    else:
                        body_parts.append(combined)

                # ── Tables ────────────────────────────────────────────────────
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            table_rows.append(" | ".join(cells))

                # ── Images ────────────────────────────────────────────────────
                try:
                    if shape.shape_type == _PICTURE_TYPE:
                        image_blobs.append(shape.image.blob)
                except Exception:
                    pass

            # ── Speaker notes ─────────────────────────────────────────────────
            notes_text = ""
            try:
                if slide.has_notes_slide:
                    nf = slide.notes_slide.notes_text_frame
                    notes_text = "\n".join(
                        "".join(run.text for run in para.runs).strip()
                        for para in nf.paragraphs
                    ).strip()
            except Exception:
                pass

            slides_data.append({
                "slide_number": slide_num,
                "title":        title,
                "body":         "\n\n".join(body_parts),
                "notes":        notes_text,
                "tables":       "[TABLE]\n" + "\n".join(table_rows) + "\n[/TABLE]"
                                if table_rows else "",
                "image_blobs":  image_blobs,
            })

        logger.info("  Extracted content from %d slides.", len(slides_data))
        return slides_data

    except Exception as exc:
        logger.error("  PPTX extraction failed: %s", exc)
        return []
    
def parse_slide_document(
    slides_data: List[Dict],
    source_key:  str,
    source_name: str,
) -> List[Dict]:
    """
    Convert extracted PPTX slide data into the standard CSV chunk format.

    One chunk per slide.  Slides with no extractable text but with images
    are OCR'd.  Oversized slides are split at sentence boundaries.
    Empty/decorative slides are skipped.
    """
    rows: List[Dict] = []
    last_title = ""   # carries section context for untitled slides

    for slide in slides_data:
        num    = slide["slide_number"]
        title  = slide["title"].strip()
        body   = slide["body"].strip()
        notes  = slide["notes"].strip()
        tables = slide["tables"].strip()

        _NOTES_PLACEHOLDER = re.compile(r"^click to add notes$", re.IGNORECASE)
        if _NOTES_PLACEHOLDER.match(notes):
            notes = ""

        # ── OCR any images on text-sparse slides ──────────────────────────────
        ocr_parts: List[str] = []
        if slide["image_blobs"] and not body:
            for img_bytes in slide["image_blobs"]:
                ocr_text = ocr_image_bytes(img_bytes, "image/png")
                if ocr_text and len(ocr_text.strip()) >= MIN_CONTENT_CHARS:
                    ocr_parts.append(f"[IMAGE TEXT]\n{ocr_text.strip()}\n[/IMAGE TEXT]")

        # ── Assemble chunk content ─────────────────────────────────────────────
        content_parts = []
        if body:
            content_parts.append(body)
        if tables:
            content_parts.append(tables)
        content_parts.extend(ocr_parts)
        if notes:
            content_parts.append(f"[SPEAKER NOTES]\n{notes}")

        content = "\n\n".join(content_parts).strip()
        if len(content) < MIN_CONTENT_CHARS:
            continue   # blank / purely decorative slide

        # ── Build section label ────────────────────────────────────────────────
        if title:
            section    = f"Slide {num}: {title}"
            last_title = title
        elif last_title:
            section    = f"Slide {num} ({last_title} – cont.)"
        else:
            section    = f"Slide {num}"

        # ── Split if over size limit ───────────────────────────────────────────
        if len(content) > MAX_GENERAL_CHUNK_CHARS:
            sub_chunks = _split_at_sentences(content, MAX_GENERAL_CHUNK_CHARS)
        else:
            sub_chunks = [content]

        for i, chunk in enumerate(sub_chunks):
            chunk = chunk.strip()
            if len(chunk) < MIN_CONTENT_CHARS:
                continue
            rows.append({
                "source_key":    source_key,
                "source_name":   source_name,
                "part":          title or last_title,
                "division":      "",
                "subdivision":   "",
                "section":       section if i == 0 else f"{section} (part {i + 1})",
                "section_title": title,
                "content":       chunk,
            })

    logger.info(
        "  Parsed %d chunk(s) from %d slide(s) in '%s'.",
        len(rows), len(slides_data), source_name,
    )
    return rows

# ─── Standalone image source ──────────────────────────────────────────────────

def process_image_source(source: SourceEntry) -> List[Dict]:
    """
    OCR a standalone image file (PNG, JPG, JPEG, WEBP) and parse the
    resulting text with parse_with_ai.

    Use this for scanned single-page documents, notices, or diagrams
    with embedded text that are delivered as image files rather than PDFs.
    """
    path = source.local_path
    if not path or not os.path.exists(path):
        logger.error("  Image not found: %s", path)
        return []

    _MIME = {
        ".png": "image/png", ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg", ".webp": "image/webp",
    }
    mime_type = _MIME.get(os.path.splitext(path)[1].lower(), "image/png")

    try:
        with open(path, "rb") as fh:
            image_bytes = fh.read()
    except Exception as exc:
        logger.error("  Failed to read image '%s': %s", path, exc)
        return []

    logger.info("  Running OCR on image (%d bytes) …", len(image_bytes))
    ocr_text = ocr_image_bytes(image_bytes, mime_type)

    if not ocr_text or len(ocr_text.strip()) < MIN_CONTENT_CHARS:
        logger.error(
            "  OCR returned insufficient text for '%s'. "
            "Check that '%s' is running via 'ollama ps'.",
            source.key, OCR_MODEL,
        )
        return []

    logger.info("  OCR extracted %d chars.", len(ocr_text))
    cleaned = clean_text(ocr_text)
    return parse_with_ai(cleaned, source.key, source.name)

# ─── CSV source handling ──────────────────────────────────────────────────────

def process_csv_source(source: SourceEntry) -> List[Dict]:
    """
    Load an external CSV and normalise it to the standard schema.

    Expected columns (flexible): Part, Division, Subdivision, Section,
    Subsection/Section_Title, Content.
    """
    import pandas as pd  # noqa: PLC0415

    path = source.local_path
    if not path or not os.path.exists(path):
        logger.error("  CSV not found: %s", path)
        return []

    try:
        df = pd.read_csv(path, encoding="utf-8", dtype=str).fillna("")
        df.columns = [c.strip() for c in df.columns]
    except Exception as exc:
        logger.error("  Failed to read CSV: %s", exc)
        return []

    # Require at least a content column
    content_col = next(
        (c for c in df.columns if c.lower() in {"content", "text", "body"}),
        None,
    )
    if not content_col:
        logger.error(
            "  CSV has no recognisable content column. Found: %s", list(df.columns)
        )
        return []

    rows: List[Dict] = []
    for _, row in df.iterrows():
        content = str(row.get(content_col, "")).strip()
        if len(content) < MIN_CONTENT_CHARS:
            continue

        rows.append({
            "source_key":    source.key,
            "source_name":   source.name,
            "part":          str(row.get("Part", row.get("part", ""))).strip(),
            "division":      str(row.get("Division", row.get("division", ""))).strip(),
            "subdivision":   str(row.get("Subdivision", row.get("subdivision", ""))).strip(),
            "section":       str(row.get("Section", row.get("section", ""))).strip(),
            "section_title": str(
                row.get("Subsection",
                row.get("Section_Title",
                row.get("section_title",
                row.get("Title", ""))))
            ).strip(),
            "content":       content,
        })

    logger.info("  Loaded %d rows from CSV '%s'.", len(rows), source.name)
    return rows

# ─── Write CSV ────────────────────────────────────────────────────────────────

def write_csv(rows: List[Dict], output_path: str) -> None:
    """Write rows to the standard processed CSV format."""
    tmp = output_path + ".tmp"
    try:
        with open(tmp, "w", newline="", encoding="utf-8") as fh:
            writer = csv.DictWriter(fh, fieldnames=CSV_COLUMNS, extrasaction="ignore")
            writer.writeheader()
            writer.writerows(rows)
        shutil.move(tmp, output_path)
        logger.info("  Saved %d rows -> %s", len(rows), output_path)
    except Exception as exc:
        logger.error("  Failed to write CSV %s : %s", output_path, exc)
        try:
            os.unlink(tmp)
        except OSError:
            pass
        raise

# ─── Main processing function ─────────────────────────────────────────────────

def process_source(source: SourceEntry, force: bool = False) -> bool:
    """
    Process one source entry: extract, clean, parse, and write CSV.

    Routing by doc_type
    -------------------
        "act"     -> parse_legal_document()     (PART/Section structure)
        "general" -> parse_general_document()   (numbered paragraphs / headings)
        "faq"     -> parse_faq_document()        (Q/A pairs)

    Returns True on success.
    """
    logger.info("")
    logger.info("=" * 60)
    logger.info(
        "Processing: %s  [%s]  doc_type=%s",
        source.name, source.key, source.doc_type,
    )
    logger.info("=" * 60)

    if source.is_processed and not force:
        logger.info("  Already processed. Use --force to re-process.")
        return True

    rows: List[Dict] = []

    if source.type == "pdf":
        pdf_bytes = get_pdf_bytes(source)
        if not pdf_bytes:
            logger.error("  Could not acquire PDF for '%s'. Skipping.", source.key)
            return False
        raw_text = extract_pdf_text(pdf_bytes)   # now OCR-aware
        if not raw_text:
            logger.error("  Text extraction failed for '%s'. Skipping.", source.key)
            return False
        cleaned = clean_text(raw_text)
        logger.info("  Parsing document (doc_type=%s) …", source.doc_type)
        if source.doc_type == "act":
            rows = parse_legal_document(cleaned, source.key, source.name)
        elif source.doc_type == "general":
            rows = parse_general_document(cleaned, source.key, source.name)
        elif source.doc_type == "faq":
            rows = parse_faq_document(cleaned, source.key, source.name)
        elif source.doc_type == "gazette":
            rows = parse_gazette_document(cleaned, source.key, source.name)
        elif source.doc_type in ("others", "slide"):
            rows = parse_with_ai(cleaned, source.key, source.name)
        else:
            logger.error("  Unknown doc_type '%s'. Skipping.", source.doc_type)
            return False

    elif source.type == "image":
        rows = process_image_source(source)

    elif source.type == "pptx":
        path = source.local_path
        if not path or not os.path.exists(path):
            logger.error("  PPTX file not found: %s", path)
            return False
        try:
            with open(path, "rb") as fh:
                pptx_bytes = fh.read()
        except Exception as exc:
            logger.error("  Failed to read PPTX '%s': %s", path, exc)
            return False
        slides_data = extract_pptx_content(pptx_bytes)
        if not slides_data:
            logger.error("  No slides extracted from '%s'.", source.key)
            return False
        rows = parse_slide_document(slides_data, source.key, source.name)

    elif source.type == "csv":
        rows = process_csv_source(source)

    else:
        logger.error("  Unknown source type '%s'. Skipping.", source.type)
        return False

    if not rows:
        logger.error(
            "  No rows produced for '%s'.\n"
            "  Check the source file and try --force to re-run.\n"
            "  If the PDF is scanned (image-only), text extraction will fail.",
            source.key,
        )
        return False

    # Step 5: write output CSV
    try:
        write_csv(rows, source.output_path)
    except Exception:
        return False
    return True


# ─── CLI ──────────────────────────────────────────────────────────────────────

def list_sources(sources: List[SourceEntry]) -> None:
    """Print a summary table of all enabled sources."""
    print(f"\n{'KEY':<40} {'CAT':<20} {'TYPE':<9} {'DONE':<5} {'PATH/URL'}")
    print("-" * 110)
    for s in sources:
        done    = "✓" if s.is_processed else "✗"
        loc     = s.local_path or s.url or "(none)"
        if loc and len(loc) > 45:
            loc = "…" + loc[-43:]
        print(f"{s.key:<40} {s.category:<20} {s.doc_type:<9} {done:<5} {loc}")
    print(f"\n  Total: {len(sources)} enabled sources\n")


def scan_sources() -> None:
    """
    Dry-run: show what files would be discovered without processing anything.
    Useful for verifying folder structure before a long preprocess run.
    """
    print(f"\nScanning '{SOURCES_DIR}' …\n")
    folder  = _discover_folder_sources()
    json_ov = _load_json_overrides()

    if not folder and not json_ov:
        print("  No sources found.")
        print(f"  Create folders like:  {SOURCES_DIR}/Legislations/act/")
        print(f"  Then drop your PDFs into the appropriate folder.\n")
        return

    print(f"{'KEY':<40} {'CATEGORY':<20} {'DOC_TYPE':<9} {'SOURCE':<8} {'FILE'}")
    print("-" * 110)
    for key, s in sorted(folder.items()):
        override = " (JSON override)" if key in json_ov else ""
        fname = os.path.basename(s.local_path or "")
        print(f"{key:<40} {s.category:<20} {s.doc_type:<9} {'folder':<8} {fname}{override}")
    for key, s in sorted(json_ov.items()):
        if key not in folder:
            loc = s.url or s.local_path or "(none)"
            print(f"{key:<40} {s.category:<20} {s.doc_type:<9} {'JSON':<8} {loc}")

    print(f"\n  {len(folder)} folder source(s), {len(json_ov)} JSON entry/entries")
    print(f"  Run without --scan to process all.\n")

# ─── AI chunking helpers ──────────────────────────────────────────────────────

def _sanitize_for_llm(text: str) -> str:
    """
    Strip extraction artifacts that confuse the AI chunker or produce invalid
    JSON — OCR markers, table wrapper tags, Unicode garbage characters.
    Called once per page-slice before the LLM request.
    """
    # Remove OCR page markers injected by extract_pdf_text()
    text = re.sub(r'\[OCR PAGE \d+\]\n?', '', text)
    # Remove table wrapper tags; keep the pipe-delimited row content
    text = re.sub(r'\[(?:/?)TABLE\]\n?', '', text)
    # Remove image-text wrapper tags from PPTX OCR
    text = re.sub(r'\[(?:/?)IMAGE TEXT\]\n?', '', text)
    # Strip Unicode private-use area and replacement characters
    text = re.sub(r'[\uE000-\uF8FF\uFFFD\uFFFE\uFFFF]', '', text)
    # Normalise to 2 consecutive newlines maximum
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _repair_json_strings(s: str) -> str:
    """
    Escape control characters that appear as literal bytes inside JSON string
    values — the most common LLM output defect.

    Standard JSON forbids unescaped U+0000-U+001F inside strings.
    The existing cleanup regex skips \\x0a (newline) and \\x0d (CR) because
    they are structurally valid outside strings; this function handles the
    case where they appear INSIDE strings.

    Uses a single-pass character state machine — O(n), no regex.
    """
    out: list = []
    in_str = False
    escaped = False
    _ESC_MAP = {'\n': '\\n', '\r': '\\r', '\t': '\\t'}

    for ch in s:
        if escaped:
            out.append(ch)
            escaped = False
        elif ch == '\\' and in_str:
            out.append(ch)
            escaped = True
        elif ch == '"':
            in_str = not in_str
            out.append(ch)
        elif in_str and ch in _ESC_MAP:
            out.append(_ESC_MAP[ch])          # \n → \\n, \r → \\r, \t → \\t
        elif in_str and ord(ch) < 0x20:
            pass                              # drop remaining C0 control chars
        else:
            out.append(ch)

    return ''.join(out)


def _split_pages_at_boundary(text: str, max_chars: int) -> list:
    """
    Split text into page-slices at paragraph boundaries, not mid-character.
    This prevents the LLM from receiving a fragment that ends mid-sentence.
    """
    if len(text) <= max_chars:
        return [text]

    pages = []
    while len(text) > max_chars:
        # Prefer paragraph break
        cut = text.rfind('\n\n', 0, max_chars)
        if cut == -1 or cut < max_chars // 2:
            # Fall back to any newline
            cut = text.rfind('\n', 0, max_chars)
        if cut == -1 or cut < max_chars // 2:
            # Fall back to word boundary
            cut = text.rfind(' ', 0, max_chars)
        if cut <= 0:
            cut = max_chars                  # hard cut only as last resort
        pages.append(text[:cut].strip())
        text = text[cut:].strip()

    if text:
        pages.append(text)
    return pages

def parse_with_ai(
    text: str,
    source_key: str,
    source_name: str,
) -> List[Dict]:
    """
    Fallback parser: uses the local Ollama LLM to chunk any document
    regardless of format. Slower (5-15s/page) but handles any layout.
    Output is the same CSV schema as all other parsers.
    """
    import json, requests

    # ── System message: suppresses thinking output + enforces JSON-only ──────
    CHUNK_SYSTEM = (
        "/no_think\n"
        "You output ONLY raw JSON. No explanations. No markdown. "
        "No preamble. No postamble. No code fences."
    )

    # ── User prompt ───────────────────────────────────────────────────────────
    CHUNK_PROMPT = """\
/no_think
Split the text below into retrieval chunks for a Malaysian legal document system.

OUTPUT RULES — STRICTLY ENFORCED
- Your ENTIRE response must be a single JSON array.
- The FIRST character must be [ and the LAST character must be ].
- No markdown fences. No explanations. No text before [ or after ].
- Inside each "content" string, represent newlines as \\n (escaped), NOT as literal line breaks.
- Minimum content length: 60 characters.

CHUNKING RULES
- Keep sub-items (a)(b)(c) together with their parent clause.
- Keep a table together with the text that introduces it.
- Do NOT split a sentence across chunks.
- For Q&A: one chunk = one question + its complete answer.

RESPONSE SCHEMA — each array element:
{{
  "part":          "PART name or section group, or empty string",
  "section":       "Section number/heading, or empty string",
  "section_title": "Short descriptive title, or empty string",
  "content":       "Full chunk text with \\n for newlines"
}}

TEXT:
{text}"""

    # ── JSON extraction with 4-pass repair ────────────────────────────────────
    def _try_parse(raw: str) -> list:
        start = raw.find('[')
        end   = raw.rfind(']')
        if start == -1 or end == -1 or start >= end:
            raise ValueError("No JSON array found in model response")
        candidate = raw[start:end + 1]

        # Pass 1 — direct parse (fast path for well-formed output)
        try:
            return json.loads(candidate)
        except json.JSONDecodeError:
            pass

        # Pass 2 — escape literal newlines/CR/tab inside string values
        try:
            return json.loads(_repair_json_strings(candidate))
        except json.JSONDecodeError:
            pass

        # Pass 3 — strip remaining C0 controls + bad backslashes
        try:
            c = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', candidate)
            c = re.sub(r'(?<!\\)\\(?!["\\/bfnrtu])', r'\\\\', c)
            return json.loads(c)
        except json.JSONDecodeError:
            pass

        # Pass 4 — combine both repairs
        fixed = _repair_json_strings(candidate)
        fixed = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', fixed)
        fixed = re.sub(r'(?<!\\)\\(?!["\\/bfnrtu])', r'\\\\', fixed)
        return json.loads(fixed)   # raise if still broken — triggers fallback

    # ── Per-page worker ───────────────────────────────────────────────────────
    def _process_page(args):
        page_num, page_text = args
        page_text = _sanitize_for_llm(page_text)
        if len(page_text.strip()) < 60:
            return page_num, []

        prompt_body = CHUNK_PROMPT.replace("{text}", page_text)
        page_rows   = []

        for attempt in range(2):   # one retry on any failure
            try:
                r = requests.post(
                    "http://localhost:11434/api/chat",
                    json={
                        "model": os.environ.get("CHATSSM_LLM_MODEL", "qwen3:1.7b"),
                        "messages": [
                            {"role": "system", "content": CHUNK_SYSTEM},
                            {"role": "user",   "content": prompt_body},
                        ],
                        "stream": False,
                        "options": {
                            "temperature": 0.0,
                            "num_predict": 3000,
                            "num_ctx":     6144,
                            "num_gpu":     CHUNK_NUM_GPU,
                        },
                    },
                    timeout=(10, 180),
                )

                # ── HTTP guard ────────────────────────────────────────────────
                if r.status_code != 200:
                    raise ValueError(
                        f"Ollama returned HTTP {r.status_code}: {r.text[:200]}"
                    )

                raw = r.json().get("message", {}).get("content", "").strip()

                # ── Strip thinking blocks (closed and unclosed) ───────────────
                raw = re.sub(r"<think>.*?</think>", "", raw, flags=re.DOTALL)
                raw = re.sub(r"<think>.*$",         "", raw, flags=re.DOTALL)
                # ── Strip markdown fences ─────────────────────────────────────
                raw = re.sub(r"```(?:json)?\s*", "", raw)
                raw = raw.strip()

                parsed = _try_parse(raw)

                for item in parsed:
                    content = str(item.get("content", "")).strip()
                    if len(content) < 60:
                        continue
                    page_rows.append({
                        "source_key":    source_key,
                        "source_name":   source_name,
                        "part":          str(item.get("part",          "")),
                        "division":      "",
                        "subdivision":   "",
                        "section":       str(item.get("section",       "")),
                        "section_title": str(item.get("section_title", "")),
                        "content":       content,
                    })
                break   # success — exit retry loop

            except Exception as exc:
                logger.warning(
                    "AI chunking attempt %d failed for page %d: %s",
                    attempt + 1, page_num, exc,
                )
                if attempt == 1:   # both attempts failed — store raw text
                    stripped = page_text.strip()
                    if len(stripped) >= 60:
                        page_rows.append({
                            "source_key":    source_key,
                            "source_name":   source_name,
                            "part":          "",
                            "division":      "",
                            "subdivision":   "",
                            "section":       f"Page {page_num + 1}",
                            "section_title": "",
                            "content":       stripped,
                        })

        return page_num, page_rows

    # ── Split and dispatch ─────────────────────────────────────────────────────
    page_size = 2500   # was 1000 — too small, produced mid-sentence fragments
    pages   = _split_pages_at_boundary(text, page_size)
    workers = int(os.environ.get("CHATSSM_CHUNK_WORKERS", "3"))

    rows_by_page: Dict[int, List[Dict]] = {}
    with ThreadPoolExecutor(max_workers=workers) as pool:
        futures = {pool.submit(_process_page, (i, p)): i for i, p in enumerate(pages)}
        for future in as_completed(futures):
            page_num, page_rows = future.result()
            rows_by_page[page_num] = page_rows

    rows: List[Dict] = []
    for i in range(len(pages)):
        rows.extend(rows_by_page.get(i, []))

    logger.info("  AI parser: produced %d chunks from '%s'.", len(rows), source_name)
    return rows

def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "ChatSSM Offline Preprocessor.\n\n"
            "Drop PDFs into knowledge_base/sources/<Category>/<doc_type>/ "
            "and run this script — no JSON editing needed."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Folder structure example:\n"
            "  knowledge_base/sources/Legislations/act/companies_act_2016.pdf\n"
            "  knowledge_base/sources/Practice Notes/general/pn_3_2018.pdf\n"
            "  knowledge_base/sources/FAQ/faq/faq_incorporation.pdf\n\n"
            "Valid categories: Legislations | Practice Notes | Practice Directives |\n"
            "                  Guidelines | Circular | FAQ | Forms\n"
            "Valid doc_types:  act | general | faq\n"
        ),
    )
    parser.add_argument("--key",   metavar="KEY",  help="Process only this source key.")
    parser.add_argument("--force", action="store_true", help="Re-process even if CSV exists.")
    parser.add_argument("--list",  action="store_true", help="List all enabled sources, then exit.")
    parser.add_argument("--scan",  action="store_true", help="Preview discovered files (dry-run), then exit.")
    args = parser.parse_args()

    if args.scan:
        scan_sources()
        return

    sources = load_sources()

    if args.list:
        list_sources(sources)
        return

    if args.key:
        targets = [s for s in sources if s.key == args.key]
        if not targets:
            logger.error("No source found with key '%s'.", args.key)
            logger.error("Run --list or --scan to see available keys.")
            sys.exit(1)
    else:
        targets = sources
        if not targets:
            logger.error(
                "No sources found.\n"
                "  Drop PDFs into knowledge_base/sources/<Category>/<doc_type>/\n"
                "  and run again, or add entries to knowledge_sources.json."
            )
            sys.exit(1)

    logger.info("Will process %d source(s).", len(targets))

    results: Dict[str, List[str]] = {"ok": [], "failed": [], "skipped": []}
    for source in targets:
        if source.is_processed and not args.force:
            results["skipped"].append(source.key)
            continue
        ok = process_source(source, force=args.force)
        (results["ok"] if ok else results["failed"]).append(source.key)

    print("\n" + "=" * 60)
    print("PREPROCESSING COMPLETE")
    print("=" * 60)
    print(f"  Processed : {len(results['ok'])}   → {results['ok']}")
    print(f"  Skipped   : {len(results['skipped'])} (already done — use --force to redo)")
    print(f"  Failed    : {len(results['failed'])}   → {results['failed']}")
    if results["ok"]:
        print(f"\n  Output files: {PROCESSED_DIR}/")
        print("  Start app:    streamlit run chatssm_app.py")
    print()

if __name__ == '__main__':
    main()